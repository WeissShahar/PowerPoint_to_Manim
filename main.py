import subprocess
import os
import webbrowser

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_AUTO_SHAPE_TYPE

# Initialize variables
presentation_path = 'presentations/PathPlanning (1).pptx'
output_path = 'demonstration_for_Roi.py'
image_dir = 'extracted_images'  # Directory for storing extracted images
slides_shapes_info = []  # Store shapes info for all slides
global_shapes = {}  # Global dictionary to store shapes info with shape_id as key
background_color = 'WHITE'  # Manim color for background

# Ensure the image directory exists
if not os.path.exists(image_dir):
    os.makedirs(image_dir)


def get_start_point_and_end_point(shape):
    """Get start and end point of a line."""
    start_point = [
        convert_margin_to_points(shape.begin_x) / 72 - frame_width / 2,
        frame_height / 2 - convert_margin_to_points(shape.begin_y) / 72,
        0
    ]
    end_point = [
        convert_margin_to_points(shape.end_x) / 72 - frame_width / 2,
        frame_height / 2 - convert_margin_to_points(shape.end_y) / 72,
        0
    ]
    return start_point, end_point


def convert_margin_to_points(margin):
    """Convert margin from Emu to points."""
    return margin / 12700  # 1 point = 12700 EMUs


# Helper Functions
def convert_position(shape, frame_width, frame_height):
    """Convert PowerPoint coordinates to Manim coordinates."""
    return [
        shape.left.pt / 72 - frame_width / 2 + shape.width.pt / 72 / 2,
        frame_height / 2 - shape.top.pt / 72 - shape.height.pt / 72 / 2,
        0
    ]


def extract_shapes_from_slide(slide, frame_width, frame_height, slide_index):
    """Function to extract relevant info from a slide and store shape details."""
    extracted_shapes = []
    for shape in slide.shapes:
        shape_info = {
            'id': shape.shape_id,
            'type': None,
            'position': convert_position(shape, frame_width, frame_height),
            'dimensions': (shape.width.pt / 72, shape.height.pt / 72),
            'text': shape.text if shape.has_text_frame and shape.text else None,
            'image_path': None
        }

        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            if shape.name.startswith('Line'):
                shape_info['type'] = 'line'
            elif shape.auto_shape_type == MSO_AUTO_SHAPE_TYPE.RECTANGLE:
                shape_info['type'] = 'rectangle'
            elif shape.auto_shape_type == MSO_AUTO_SHAPE_TYPE.OVAL:
                shape_info['type'] = 'oval'

            else:
                continue
        elif "Arrow" in shape.name:  # Detect arrows
            shape_info['type'] = 'arrow'
            shape_info['width'] = shape.width / 914440 if shape.width / 914440 >= 0.1 else 1
            start_point, end_point = get_start_point_and_end_point(shape)
            shape_info['dimensions'] = (start_point, end_point)
        elif "Connector" in shape.name:
            shape_info['type'] = 'line'
            start_point, end_point = get_start_point_and_end_point(shape)
            if shape.line.dash_style:
                shape_info['dash_style'] = 'dashed'
            else:
                shape_info['dash_style'] = 'solid'

                # Handle the line color if it's defined
            color = shape.line.color
            if color and hasattr(color, 'rgb') and color.rgb:
                shape_info['color'] = color.rgb
            else:
                shape_info['color'] = '000000'  # Default color if not specified
            shape_info['dimensions'] = (start_point, end_point)

        elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX or shape.has_text_frame:  # Text box
            if shape.has_text_frame and shape.text:
                shape_info['type'] = 'text'
            else:
                continue

        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:  # Image
            shape_info['type'] = 'image'
            image = shape.image
            image_bytes = image.blob
            image_filename = os.path.join(image_dir, f"slide_{slide_index}_image_{shape.shape_id}.png")
            with open(image_filename, 'wb') as image_file:
                image_file.write(image_bytes)
            shape_info['image_path'] = image_filename

        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:  # Table
            shape_info['type'] = 'table'
            table = shape.table
            table_data = []
            font_sizes = []
            for row in table.rows:
                row_data = []
                row_font_sizes = []
                for cell in row.cells:
                    row_data.append(cell.text)
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            row_font_sizes.append(run.font.size.pt)
                table_data.append(row_data)
                font_sizes.append(row_font_sizes)

            shape_info['font_sizes'] = font_sizes
            shape_info['table_data'] = table_data

        if shape.shape_type != MSO_SHAPE_TYPE.TEXT_BOX and shape.has_text_frame and shape.text and shape.shape_type:
            text_info = {
                'id': f"{shape.shape_id}_text",
                'type': 'text',
                'text': shape.text,
                'position': convert_position(shape, frame_width, frame_height),  # Align text to the center of the shape
                'dimensions': (shape.width.pt / 72, shape.height.pt / 72),
                'image_path': None
            }
            extracted_shapes.append(text_info)  # Add text as a separate shape

        extracted_shapes.append(shape_info)
        global_shapes[shape_info['id']] = shape_info

    return extracted_shapes


def generate_manim_code(slides_shapes_info, background_color, frame_width, frame_height):
    """Generate Manim code as a string."""
    manim_code = f"""
from manim import *
from manim_slides import Slide

class GeneratedPresentation(Slide):
    def construct(self):
        self.camera.background_color = {background_color}
        config.frame_width = {frame_width}
        config.frame_height = {frame_height}
    """
    for i, slide_shapes_info in enumerate(slides_shapes_info):
        slide_code = f"\n        # Slide {i + 1}\n        self.clear()\n"

        # Add slide number as a text mobject in the bottom-right corner
        slide_number_position = [frame_width / 2 - 1, -frame_height / 2 + 0.5, 0]  # Adjust as necessary
        slide_code += f"        slide_number = Text('Slide {i + 1}', font_size=18, color=BLACK).move_to({slide_number_position})\n"
        slide_code += "        self.add(slide_number)\n"
        for shape_info in slide_shapes_info:
            if shape_info['type'] == 'rectangle':
                slide_code += f"        mobject = Rectangle(width={shape_info['dimensions'][0]}, height={shape_info['dimensions'][1]}, color=BLACK)\n"
            elif shape_info['type'] == 'oval':
                slide_code += f"        mobject = Ellipse(width={shape_info['dimensions'][0]}, height={shape_info['dimensions'][1]}, color=BLACK)\n"
            elif shape_info['type'] == 'text':
                slide_code += f"        mobject = Text('''{shape_info['text']}''', font_size=24, color=BLACK)\n"
            elif shape_info['type'] == 'image':
                slide_code += f"        mobject = ImageMobject('{shape_info['image_path']}')\n"
                slide_code += f"        mobject.width, mobject.height = {shape_info['dimensions']}\n"
            elif shape_info['type'] == 'line':
                    start_point, end_point = shape_info['dimensions']
                    color_hex = f"0x{shape_info['color']}"
                    if shape_info.get('dash_style') == 'dashed':
                        slide_code += f"        mobject = DashedLine(start={start_point}, end={end_point}, color=ManimColor.from_rgb({color_hex}))\n"
                    else:
                        slide_code += f"        mobject = Line(start={start_point}, end={end_point}, color=ManimColor.from_rgb({color_hex}))\n"

            elif shape_info['type'] == 'arrow':
                start_point, end_point = shape_info['dimensions']
                slide_code += f"        mobject = Arrow(start={start_point}, end={end_point}, color=BLACK, buff=1, max_tip_length_to_length_ratio=0.1, stroke_width = {shape_info['width']}  )\n"
            elif shape_info['type'] == 'table':
                table_data = shape_info['table_data']
                line_config = shape_info.get('line_config', {"stroke_color": "BLACK", "stroke_width": 2})
                element_to_mobject_config = shape_info.get('element_to_mobject_config', {"color": "BLACK"})

                table_data_str = str(table_data).replace("'", '"')  # Convert to a string representation
                line_config_str = str(line_config).replace("'", '"')
                element_to_mobject_config_str = str(element_to_mobject_config).replace("'", '"')

                slide_code += f"""
        table_data = {table_data_str}
        mobject = MathTable(
            table_data,
            include_outer_lines=True,
            line_config={line_config_str},
            element_to_mobject_config={element_to_mobject_config_str},
        )
"""
                slide_code += "        mobject.scale(0.4)\n"  # Adjust scale if necessary

            slide_code += f"        mobject.move_to({shape_info['position']})\n"
            slide_code += "        self.add(mobject)\n"

        slide_code += "        self.wait(1)\n        self.next_slide()\n"
        manim_code += slide_code

    manim_code += "\n"

    return manim_code


# Main Execution
if __name__ == "__main__":
    presentation = Presentation(presentation_path)

    # Define frame width and height based on the PowerPoint slide dimensions
    frame_width = presentation.slide_width.pt / 72
    frame_height = presentation.slide_height.pt / 72

    slide = presentation.slides[5]
    shapes = extract_shapes_from_slide(slide, frame_width, frame_height, 1)
    slides_shapes_info.append(shapes)

    # for slide_index, slide in enumerate(presentation.slides):
    #     shapes = extract_shapes_from_slide(slide, frame_width, frame_height, slide_index)
    #     slides_shapes_info.append(shapes)

    manim_code = generate_manim_code(slides_shapes_info, background_color, frame_width, frame_height)

    with open(output_path, 'w', encoding='utf-8') as manim_script:
        manim_script.write(manim_code)

    print(f"Manim code generated and saved to {output_path}")

    # Automate rendering and conversion
    print("Running Manim rendering...")
    subprocess.run(['manim', '-ql', output_path, 'GeneratedPresentation'])

    print("Converting to HTML...")
    subprocess.run(['manim-slides', 'convert', 'GeneratedPresentation', 'pre.html'])

    print("Process completed. HTML presentation saved as pre.html.")

    # Open the HTML presentation in the default web browser
    html_path = os.path.abspath('pre.html')
    print(f"Opening HTML presentation: {html_path}")
    webbrowser.open(f'file://{html_path}')
